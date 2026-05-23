"""
slide_generator.py — Professional PPTX builder for EduAI Classroom.
Supports 7 bullet points per slide with a clean dark-academic theme.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0C, 0x15, 0x20)
C_BG_MID    = RGBColor(0x16, 0x25, 0x40)
C_BG_LIGHT  = RGBColor(0xF4, 0xF7, 0xFC)
C_ACCENT    = RGBColor(0x3B, 0xC4, 0x7A)
C_BLUE      = RGBColor(0x5B, 0xA4, 0xF5)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT = RGBColor(0x1A, 0x25, 0x35)
C_MUTED     = RGBColor(0x6A, 0x7C, 0x95)
C_ROW_A     = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW_B     = RGBColor(0xED, 0xF2, 0xFB)
C_HDR       = RGBColor(0x0C, 0x15, 0x20)

W = Inches(13.33)
H = Inches(7.5)


def _rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def _text(slide, text, left, top, width, height,
          size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = str(text)
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = "Calibri"
    return tb

def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


# ── Title slide ───────────────────────────────────────────────────────────────
def build_title_slide(prs, slide_data, topic_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(slide, C_BG_DARK)

    # Gradient-like layered rectangles for depth
    _rect(slide, 0, 0, W, H, C_BG_DARK)
    _rect(slide, Inches(7), Inches(-1), Inches(7), Inches(5.5), C_BG_MID)

    # Left accent bar
    _rect(slide, 0, 0, Inches(0.18), H, C_ACCENT)

    # Decorative circle top-right
    circ = slide.shapes.add_shape(9, W - Inches(4), Inches(-1.5), Inches(5), Inches(5))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x1A, 0x38, 0x58)
    circ.line.fill.background()

    # Small label
    _text(slide, "🎓  AI CLASSROOM · EDUAI",
          Inches(0.45), Inches(1.2), Inches(8), Inches(0.5),
          size=12, bold=True, color=C_ACCENT)

    # Main title — large
    title = slide_data.get("title", topic_title)
    _text(slide, title,
          Inches(0.45), Inches(2.0), Inches(8.6), Inches(2.8),
          size=44, bold=True, color=C_WHITE)

    # Subtitle from first bullet
    bullets = slide_data.get("bullets", [])
    if bullets:
        _text(slide, bullets[0],
              Inches(0.45), Inches(4.55), Inches(7.5), Inches(0.8),
              size=18, color=RGBColor(0xA8, 0xC8, 0xFF), italic=True)

    # Bottom bar
    _rect(slide, 0, H - Inches(0.5), W, Inches(0.5), RGBColor(0x06, 0x0E, 0x18))
    _text(slide, "Powered by Lemonade Server  ·  Local AI  ·  EduAI",
          Inches(0.45), H - Inches(0.48), Inches(10), Inches(0.42),
          size=10, color=C_MUTED)


# ── Content slide — supports 7 bullets ───────────────────────────────────────
def build_content_slide(prs, slide_data, slide_num, total_content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_BG_LIGHT)

    # Header bar
    _rect(slide, 0, 0, W, Inches(1.15), C_HDR)
    # Accent left stripe in header
    _rect(slide, 0, 0, Inches(0.1), Inches(1.15), C_ACCENT)

    title = slide_data.get("title", f"Slide {slide_num}")
    _text(slide, title,
          Inches(0.28), Inches(0.15), W - Inches(2.4), Inches(0.85),
          size=26, bold=True, color=C_WHITE)

    # Slide counter badge
    badge_w = Inches(1.4)
    _rect(slide, W - badge_w - Inches(0.2), Inches(0.28),
          badge_w, Inches(0.55), RGBColor(0x1E, 0x2D, 0x42))
    _text(slide, f"{slide_num} / {total_content}",
          W - badge_w - Inches(0.2), Inches(0.26),
          badge_w, Inches(0.6),
          size=13, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # 7 bullet rows — calculate even vertical distribution
    bullets   = (slide_data.get("bullets") or [])[:7]
    while len(bullets) < 7:
        bullets.append("Additional insight related to this topic.")

    body_top  = Inches(1.22)
    body_bot  = H - Inches(0.42)
    body_h    = body_bot - body_top
    row_h     = body_h / 7
    gap       = Inches(0.04)

    for i, bullet in enumerate(bullets):
        y        = body_top + i * row_h
        row_bg   = C_ROW_A if i % 2 == 0 else C_ROW_B

        # Row background
        _rect(slide, Inches(0.22), y + gap/2, W - Inches(0.44), row_h - gap, row_bg)

        # Number circle
        num_sz = min(row_h - gap - Inches(0.06), Inches(0.38))
        nx     = Inches(0.30)
        ny     = y + (row_h - num_sz) / 2
        circ   = slide.shapes.add_shape(9, nx, ny, num_sz, num_sz)
        circ.fill.solid(); circ.fill.fore_color.rgb = C_ACCENT
        circ.line.fill.background()
        _text(slide, str(i + 1), nx, ny, num_sz, num_sz,
              size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Bullet text
        txt_x = Inches(0.82)
        _text(slide, bullet,
              txt_x, y + gap / 2, W - txt_x - Inches(0.3), row_h - gap,
              size=13, color=C_DARK_TEXT)

    # Footer
    _rect(slide, 0, H - Inches(0.38), W, Inches(0.38), C_HDR)
    _text(slide, "EduAI Classroom  ·  Lemonade Server",
          Inches(0.3), H - Inches(0.36), Inches(6), Inches(0.33),
          size=9, color=C_MUTED)


# ── Summary slide ─────────────────────────────────────────────────────────────
def build_summary_slide(prs, slide_data, topic_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_BG_DARK)

    _rect(slide, 0, 0, W, Inches(0.12), C_ACCENT)    # top accent bar
    _rect(slide, 0, H - Inches(0.5), W, Inches(0.5), RGBColor(0x06, 0x0E, 0x18))

    # Deco circle bottom-right
    circ = slide.shapes.add_shape(9, W - Inches(3), H - Inches(3), Inches(4), Inches(4))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x44)
    circ.line.fill.background()

    # Tag
    _text(slide, "KEY TAKEAWAYS",
          Inches(0.5), Inches(0.25), Inches(5), Inches(0.5),
          size=11, bold=True, color=C_ACCENT)

    title = slide_data.get("title", "Summary")
    _text(slide, title,
          Inches(0.5), Inches(0.72), W - Inches(1), Inches(1.1),
          size=34, bold=True, color=C_WHITE)

    bullets  = (slide_data.get("bullets") or [])[:7]
    while len(bullets) < 7:
        bullets.append("Continue exploring this important subject.")

    body_top = Inches(1.85)
    body_h   = H - body_top - Inches(0.6)
    row_h    = body_h / 7

    for i, bullet in enumerate(bullets):
        y = body_top + i * row_h
        # dot
        dot_sz = Inches(0.14)
        d = slide.shapes.add_shape(9,
            Inches(0.5), y + (row_h - dot_sz) / 2, dot_sz, dot_sz)
        d.fill.solid()
        d.fill.fore_color.rgb = C_ACCENT if i % 2 == 0 else C_BLUE
        d.line.fill.background()

        _text(slide, bullet,
              Inches(0.80), y, W - Inches(1.1), row_h,
              size=14, color=RGBColor(0xCC, 0xDD, 0xF5))

    _text(slide, "Questions? Ask the AI lecturer anytime!",
          Inches(0.3), H - Inches(0.47), W - Inches(0.6), Inches(0.38),
          size=11, color=C_BLUE, align=PP_ALIGN.CENTER, italic=True)


# ── Entry point ───────────────────────────────────────────────────────────────
def generate_presentation(content: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slides      = content.get("slides", [])
    topic_title = content.get("title", "Presentation")
    total       = len(slides)

    for idx, slide_data in enumerate(slides):
        if idx == 0:
            build_title_slide(prs, slide_data, topic_title)
        elif idx == total - 1:
            build_summary_slide(prs, slide_data, topic_title)
        else:
            build_content_slide(prs, slide_data, idx, total - 2)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    sample = {
        "title": "Introduction to Machine Learning",
        "slides": [
            {"title": "Introduction to Machine Learning",
             "bullets": ["Machine learning is a branch of AI focused on learning from data",
                         "It enables computers to improve performance without explicit programming",
                         "Developed in the 1950s, it has grown exponentially since 2010",
                         "Three main paradigms: supervised, unsupervised, reinforcement learning",
                         "Used in image recognition, NLP, recommendation systems, and more",
                         "Requires large datasets, compute power, and domain knowledge",
                         "Transforming industries from healthcare to autonomous vehicles"],
             "speaker_note": "Welcome everyone. Machine learning has fundamentally changed how computers solve problems. Rather than following explicit rules, ML models learn patterns from experience. Today we'll build a comprehensive understanding of this powerful field."},
            {"title": "Core Concepts of Machine Learning",
             "bullets": ["A model is a mathematical function mapping inputs to outputs",
                         "Training data is the labelled dataset used to teach the model",
                         "Loss function measures how far predictions are from true values",
                         "Gradient descent iteratively minimises the loss function",
                         "Overfitting occurs when a model memorises training data",
                         "Regularisation techniques prevent overfitting and improve generalisation",
                         "Hyperparameters control the learning process and model architecture"],
             "speaker_note": "At the heart of every ML system are a few core concepts. Understanding loss functions and gradient descent gives you insight into how models actually improve over training iterations."},
            {"title": "Summary: Key Takeaways",
             "bullets": ["Machine learning learns patterns from data rather than explicit rules",
                         "Supervised learning requires labelled data for training",
                         "Model evaluation on unseen test data is critical for validity",
                         "Overfitting is controlled through regularisation and validation",
                         "Deep learning uses neural networks with many layers",
                         "ML applications span healthcare, finance, robotics, and NLP",
                         "Ethical considerations are essential in any ML deployment"],
             "speaker_note": "We've covered the fundamentals of machine learning today. These seven takeaways capture the most important concepts. I encourage you to explore each area in depth through the recommended resources."},
        ]
    }
    path = generate_presentation(sample, "/tmp/test_edu.pptx")
    print(f"Generated: {path}")
